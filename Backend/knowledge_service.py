import os
import hashlib
import json
import logging
from datetime import datetime
from google import genai as _genai

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

from database import knowledge_base_collection, knowledge_chunks_collection
from chunking_service import chunk_text
from embedding_service import generate_document_embeddings
from config import GEMINI_MODEL

logger = logging.getLogger(__name__)

_GEMINI_KEY = os.getenv("GEMINI_API_KEY", "").strip('"').strip("'")
_knowledge_client = _genai.Client(api_key=_GEMINI_KEY)

def _call_gemini(prompt):
    resp = _knowledge_client.models.generate_content(model=GEMINI_MODEL, contents=prompt)
    return resp.text

def extract_file_text(path):
    """Extracts raw text from PDF, DOCX, or PPTX files."""
    ext = os.path.splitext(path)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            reader = PdfReader(path)
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
        elif ext == ".docx":
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext == ".pptx":
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        logger.error(f"File parse error during knowledge extraction: {e}")
    return text.strip()

def generate_file_hash(user_id, content):
    """Generate SHA256 hash to prevent duplicate document insertions per user."""
    hash_input = f"{user_id}:{content}"
    return hashlib.sha256(hash_input.encode('utf-8')).hexdigest()

def generate_summary_and_tags(text):
    """Uses Gemini to generate a summary, tags, and a knowledge score."""
    prompt = f"""
You are the Aura Knowledge Analyzer. Analyze the following document text and extract metadata.

Output ONLY a valid JSON object with this exact schema:
{{
  "summary": "string (A concise 1-3 sentence summary of the document's core content)",
  "tags": ["string", "string"] (3-7 highly relevant keywords/tags in lowercase),
  "knowledge_score": integer (1-10, rating how useful, dense, or important this document is)
}}

Document Text (first 30000 characters):
{text[:30000]}
"""
    try:
        resp_text = _call_gemini(prompt).strip()
        
        # Clean markdown wrappers if present
        if resp_text.startswith("```json"):
            resp_text = resp_text[7:-3].strip()
        elif resp_text.startswith("```"):
            resp_text = resp_text[3:-3].strip()
            
        return json.loads(resp_text)
    except Exception as e:
        logger.error(f"Failed to generate summary and tags: {e}")
        return {
            "summary": "Document content successfully extracted.",
            "tags": ["uncategorized"],
            "knowledge_score": 5
        }

def store_document(document_data):
    """Stores the document in MongoDB Knowledge Vault."""
    try:
        result = knowledge_base_collection.insert_one(document_data)
        return str(result.inserted_id)
    except Exception as e:
        logger.error(f"Failed to store document: {e}")
        return None

def process_document(user_id, file_path, original_filename):
    """Orchestrates the ingestion pipeline for a new document."""
    logger.info(f"📚 knowledge_service: Starting ingestion for '{original_filename}' (User: {user_id})")
    # 1. Extract Text
    content = extract_file_text(file_path)
    if not content:
        logger.error(f"📚 knowledge_service: Failed to extract text from '{original_filename}'")
        raise ValueError("Could not extract any text from the document.")

    logger.info(f"📚 knowledge_service: Extracted {len(content)} characters. Checking for duplicates...")
    # 2. Duplicate Detection
    file_hash = generate_file_hash(user_id, content)
    existing_doc = knowledge_base_collection.find_one({
        "user_id": user_id, 
        "file_hash": file_hash
    })
    
    if existing_doc:
        logger.info(f"⏭️ Skipped duplicate document: {original_filename}")
        return {"status": "skipped", "message": "Duplicate document exists", "id": str(existing_doc["_id"])}

    logger.info("📚 knowledge_service: Generating AI metadata (Summary & Tags)...")
    # 3. Generate AI Metadata
    ai_metadata = generate_summary_and_tags(content)
    logger.info(f"📚 knowledge_service: AI Metadata generated. Summary: {ai_metadata.get('summary', '')[:50]}...")

    # 4. Construct Schema
    ext = os.path.splitext(original_filename)[1].lower().replace(".", "")
    document_data = {
        "user_id": user_id,
        "document_name": original_filename,
        "document_type": ext or "unknown",
        "file_hash": file_hash,
        "content": content,
        "summary": ai_metadata.get("summary", ""),
        "tags": ai_metadata.get("tags", []),
        "word_count": len(content.split()),
        "knowledge_score": ai_metadata.get("knowledge_score", 5),
        "chunk_id": None,  # Future-proofing for Phase 4 RAG
        "uploaded_at": datetime.utcnow(),
        "last_accessed": datetime.utcnow()
    }

    logger.info("📚 knowledge_service: Storing document in knowledge_base_collection...")
    # 5. Store in MongoDB
    doc_id = store_document(document_data)
    
    if doc_id:
        logger.info(f"📚 knowledge_service: Document stored with ID {doc_id}. Proceeding to chunking & embedding...")
        # 6. Chunking & Embeddings
        chunks = chunk_text(content, chunk_size=400, overlap=50)
        chunks_with_embeddings = generate_document_embeddings(chunks)
        
        chunk_documents = []
        for c in chunks_with_embeddings:
            chunk_documents.append({
                "user_id": user_id,
                "document_id": doc_id,
                "document_name": original_filename,
                "chunk_id": c["chunk_id"],
                "chunk_text": c["chunk_text"],
                "embedding": c.get("embedding", []),
                "chunk_index": c["chunk_index"],
                "created_at": datetime.utcnow()
            })
            
        if chunk_documents:
            knowledge_chunks_collection.insert_many(chunk_documents)
            logger.info(f"✅ Stored {len(chunk_documents)} vector chunks for {original_filename}")
    
    return {
        "status": "success",
        "id": doc_id,
        "summary": document_data["summary"],
        "tags": document_data["tags"]
    }

def search_documents(user_id, query):
    """Executes a MongoDB Text Search on the knowledge base."""
    try:
        # $text search requires the text index we created in database.py
        cursor = knowledge_base_collection.find(
            {
                "user_id": user_id,
                "$text": {"$search": query}
            },
            # Projection to get the text match score
            {"score": {"$meta": "textScore"}}
        ).sort([("score", {"$meta": "textScore"})]).limit(3)
        
        return list(cursor)
    except Exception as e:
        logger.error(f"Knowledge Search Error: {e}")
        return []

def retrieve_relevant_knowledge(user_id, question):
    """Searches and formats the top relevant documents for the Gemini prompt."""
    if not question:
        return ""
        
    docs = search_documents(user_id, question)
    if not docs:
        return ""
        
    # Update last_accessed timestamp
    doc_ids = [d["_id"] for d in docs]
    knowledge_base_collection.update_many(
        {"_id": {"$in": doc_ids}},
        {"$set": {"last_accessed": datetime.utcnow()}}
    )
    
    # Format knowledge context string (Injecting ONLY name, summary, and tags as requested)
    formatted_knowledge = "Relevant Knowledge:\n"
    for doc in docs:
        formatted_knowledge += f"* Document Name: {doc.get('document_name')}\n"
        formatted_knowledge += f"  Summary: {doc.get('summary')}\n"
        tags = ", ".join(doc.get("tags", []))
        formatted_knowledge += f"  Tags: [{tags}]\n\n"
        
    return formatted_knowledge.strip()
